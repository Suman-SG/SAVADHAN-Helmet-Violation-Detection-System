from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=(0, 0, 0)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = RGBColor(*color)
    p.space_after = Pt(6)
    return txBox


def remake_poster(src_path=None, out_path='Helmet_Detection_Poster_Simple_remade.pptx'):
    # create presentation sized for A0 (33.11 x 46.81 inches)
    prs = Presentation()
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(46.81)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Palette (light, playful)
    title_color = (45, 127, 184)  # blue
    accent_color = (127, 184, 45)  # green
    warm_color = (179, 138, 45)  # warm accent
    text_color = (40, 40, 40)

    margin = Inches(0.5)

    # Title block
    add_textbox(slide, margin, Inches(0.5), prs.slide_width - 2 * margin, Inches(1.8),
                'Automated Helmet Detection System with FASTag-based Fine Deduction', font_size=48, bold=True,
                color=title_color)

    # Subtitle / authors
    add_textbox(slide, margin, Inches(2.4), Inches(20), Inches(0.6),
                'Computer Science & Engineering — Project by Team', font_size=28, bold=False, color=text_color)

    # Left column (text content)
    col_left = margin
    col_top = Inches(3.2)
    col_width = Inches(18)
    gap = Inches(0.2)

    # Sections: Abstract, Objectives, Methodology, Features, Results, Conclusion, Future Scope, Impact
    abstract = (
        'Abstract:\nAn AI-powered helmet detection and violation monitoring system using YOLOv8 and CLAHE enhancement. '
        'It performs real-time helmet detection, improves low-light visibility, preserves evidence (license plate), '
        'and reduces duplicate alerts via a 12-hour cooldown mechanism.'
    )
    add_textbox(slide, col_left, col_top, col_width, Inches(2.6), abstract, font_size=18, color=text_color)

    objectives = (
        'Objectives:\n- Detect helmeted and non-helmeted riders in real time\n'
        '- Improve detection in low-light using CLAHE\n- Reduce duplicate alerts (12-hour cooldown)\n- Store violation evidence with plate tracking (FASTag-like)'
    )
    add_textbox(slide, col_left, col_top + Inches(2.9), col_width, Inches(1.8), objectives, font_size=18, color=text_color)

    methodology = (
        'Methodology:\n'
        '1) Capture frames from CCTV/webcam; preprocess with CLAHE + gamma when low-light detected.\n'
        '2) Run YOLOv8 rider-and-helmet detection (two-stage) to classify riders and helmets.\n'
        '3) If helmet absent, crop plate region and run OCR (EasyOCR + Tesseract fallback).\n'
        '4) Use 12-hour cooldown keyed by plate to prevent repeated alerts.\n'
        '5) Store annotated evidence and metadata in DB; allow manual review and challan issuance.'
    )
    add_textbox(slide, col_left, col_top + Inches(4.9), col_width, Inches(3.4), methodology, font_size=16, color=text_color)

    features = (
        'Key Features:\n- Real-time helmet/no-helmet detection\n- CLAHE-based low-light enhancement\n- Dual-model detector for robustness\n- OCR with SR fallback for plate extraction\n- 12-hour duplicate suppression\n- Lightweight deployment (Flask + SQLite/MySQL)'
    )
    add_textbox(slide, col_left, col_top + Inches(8.4), col_width, Inches(2.2), features, font_size=16, color=text_color)

    results = (
        'Selected Results:\n- Helmet Detection Accuracy: 100% (sample)\n- Plate Detection Accuracy: 60%\n- OCR Character Accuracy: 59.6%\n- Precision: 89% | Recall: 88%\n- Inference: ~40 ms/image (model dependent)'
    )
    add_textbox(slide, col_left, col_top + Inches(10.7), col_width, Inches(1.8), results, font_size=16, color=text_color)

    conclusion = (
        'Conclusion:\nThe system improves detection in challenging lighting and reduces manual monitoring. '
        'Integration with FastTag-like tracking enables practical deployment for smart city traffic enforcement.'
    )
    add_textbox(slide, col_left, col_top + Inches(12.7), col_width, Inches(1.2), conclusion, font_size=16, color=text_color)

    future = (
        'Future Work:\n- Automatic challan generation and payment integration\n- Cloud dashboard & analytics\n- Mobile app for enforcement officers\n- Wider dataset and model fine-tuning for higher OCR accuracy'
    )
    add_textbox(slide, col_left, col_top + Inches(14.0), col_width, Inches(1.8), future, font_size=16, color=text_color)

    impact = (
        'Impact:\n- Improves road safety\n- Reduces manual monitoring effort\n- Supports smart city infrastructure and compliance campaigns'
    )
    add_textbox(slide, col_left, col_top + Inches(15.9), col_width, Inches(1.0), impact, font_size=16, color=text_color)

    # Right column: placeholders for images and QR
    right_left = col_left + col_width + gap
    right_width = prs.slide_width - right_left - margin
    # Large image placeholder
    img_box = slide.shapes.add_shape(1, right_left, col_top, right_width, Inches(10))
    # small QR placeholder
    qr_top = col_top + Inches(10.4)
    qr_box = slide.shapes.add_shape(1, right_left, qr_top, Inches(3), Inches(3))

    # Add labels for placeholders
    add_textbox(slide, right_left, col_top - Inches(0.3), right_width, Inches(0.4), 'Image Area (paste high-res images here)', font_size=14, color=warm_color)
    add_textbox(slide, right_left, qr_top + Inches(3.1), Inches(3), Inches(0.6), 'QR / Link (place QR image here)', font_size=12, color=title_color)

    # Footer with contact/github/demo links placeholder
    footer_text = ('Scan for: Live Demo | GitHub Repo | Project Report | Demo Video | Website')
    add_textbox(slide, margin, prs.slide_height - Inches(1.0), prs.slide_width - 2 * margin, Inches(0.8), footer_text, font_size=14, color=RGBColor(80, 80, 80))

    # Save output
    out_full = os.path.join(os.getcwd(), out_path)
    prs.save(out_full)
    print(f'Saved remade poster: {out_full}')


if __name__ == '__main__':
    remake_poster()
